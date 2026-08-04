import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import sqlite3
import csv
import os
import sys
import time
import io
import warnings
import shutil
import threading
import smtplib
import re
import json
from datetime import datetime
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from email.mime.image import MIMEImage
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
import numpy as np

# --- UNIVERSAL DIRECTORY FIX ---
if getattr(sys, 'frozen', False):
    if sys.platform == 'darwin' and '.app' in sys.executable:
        app_dir = os.path.dirname(os.path.dirname(os.path.dirname(sys.executable)))
        os.chdir(os.path.dirname(app_dir))
    else:
        exe_dir = os.path.dirname(sys.executable)
        if os.path.basename(exe_dir).lower() == "systemfiles":
            os.chdir(os.path.dirname(exe_dir))
        else:
            os.chdir(exe_dir)
else:
    script_dir = os.path.dirname(os.path.abspath(__file__))
    if os.path.basename(script_dir).lower() == "systemfiles":
        os.chdir(os.path.dirname(script_dir))
    else:
        os.chdir(script_dir)

warnings.filterwarnings("ignore", message="Tight layout not applied")

try:
    from PIL import ImageGrab, Image
    HAS_PILLOW = True
except ImportError:
    HAS_PILLOW = False

try:
    import requests
    HAS_REQUESTS = True
except ImportError:
    HAS_REQUESTS = False

try:
    from tkcalendar import DateEntry
    HAS_TKCALENDAR = True
except ImportError:
    HAS_TKCALENDAR = False

try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import comtypes.client
    HAS_COMTYPES = True
except ImportError:
    HAS_COMTYPES = False

WHO_GUIDELINES = {
    'M': {
        'BMI': {5: 15.2, 6: 15.3, 7: 15.5, 8: 15.7, 9: 16.1, 10: 16.7, 11: 17.3, 12: 18.0, 13: 18.8, 14: 19.6, 15: 20.4, 16: 21.1, 17: 21.8, 18: 22.4},
        'Height': {5: 110.0, 6: 116.0, 7: 121.9, 8: 127.3, 9: 132.6, 10: 137.8, 11: 143.1, 12: 149.1, 13: 156.0, 14: 163.2, 15: 170.0, 16: 173.5, 17: 175.2, 18: 176.1},
        'Weight': {5: 18.3, 6: 20.5, 7: 22.9, 8: 25.4, 9: 28.1, 10: 31.2, 11: 35.3, 12: 39.8, 13: 45.3, 14: 50.8, 15: 56.0, 16: 60.8, 17: 64.6, 18: 66.9}
    },
    'F': {
        'BMI': {5: 15.3, 6: 15.3, 7: 15.4, 8: 15.7, 9: 16.1, 10: 16.6, 11: 17.2, 12: 18.0, 13: 18.8, 14: 19.6, 15: 20.2, 16: 20.7, 17: 21.1, 18: 21.4},
        'Height': {5: 109.4, 6: 115.1, 7: 120.8, 8: 126.6, 9: 132.2, 10: 138.6, 11: 145.0, 12: 151.2, 13: 156.4, 14: 159.8, 15: 161.7, 16: 162.5, 17: 162.9, 18: 163.1},
        'Weight': {5: 18.2, 6: 20.2, 7: 22.4, 8: 25.0, 9: 28.2, 10: 31.9, 11: 36.9, 12: 41.5, 13: 45.8, 14: 49.8, 15: 53.0, 16: 55.4, 17: 56.7, 18: 57.3}
    }
}

def categorize_metric(metric, val, age, gender):
    if gender not in ['M', 'F'] or age is None or val is None: return "Within Range" 
    safe_age = max(5, min(18, age))
    try: median = WHO_GUIDELINES[gender][metric][safe_age]
    except KeyError: return "Within Range" 
    if metric == "Height": low_mult, high_mult = 0.93, 1.07
    else: low_mult, high_mult = 0.85, 1.15
    lower_bound, upper_bound = median * low_mult, median * high_mult
    if lower_bound <= val <= upper_bound: return "Within Range"
    elif (lower_bound - 5) <= val <= (upper_bound + 5): return "Borderline"
    else: return "Needs Attention"

class LVLXCommandCenter:
    def __init__(self, root):
        self.root = root
        self.all_student_data = []
        self.current_log_file = None 
        self.active_text_widget = None 
        self.review_cards = [] 
        self.current_erp = None 
        self._autosave_timer = None
        
        # --- DYNAMIC ROOT DIRECTORY FIX ---
        passed_db_path = None
        for arg in sys.argv[1:]:
            if arg.endswith('.db'):
                passed_db_path = os.path.abspath(arg)
                break
                
        if not passed_db_path:
            self.root.withdraw() 
            messagebox.showerror("Access Denied", "Direct access is disabled.\n\nPlease launch this tool through the LVLX School Portal.")
            os._exit(0)

        abs_db_path = os.path.abspath(passed_db_path)
        parent_dir = os.path.dirname(abs_db_path)
        parent_name = os.path.basename(parent_dir)
        
        if parent_name == "database":
            root_dir = os.path.dirname(os.path.dirname(os.path.dirname(parent_dir)))
            os.chdir(root_dir)
        elif parent_name == "Databases":
            root_dir = os.path.dirname(parent_dir)
            os.chdir(root_dir)
        # ----------------------------------

        raw_name = os.path.basename(passed_db_path).replace('.db', '')
        self.display_school_name = raw_name.replace('_', ' ').replace('-', ' ')

        base_reports_dir = os.path.join(os.getcwd(), "Reports")
        school_dir = os.path.join(base_reports_dir, raw_name)
        db_folder = os.path.join(school_dir, "database")
        os.makedirs(db_folder, exist_ok=True)
        
        self.db_name = os.path.join(db_folder, f"{raw_name}.db")
        
        if os.path.abspath(passed_db_path) != os.path.abspath(self.db_name):
            if os.path.exists(passed_db_path):
                shutil.move(passed_db_path, self.db_name)

        self.root.title(f"LVLX Command Center - {self.display_school_name}")
        self.root.geometry("1100x750")
        
        # --- CROSS-PLATFORM THEME ENGINE ---
        self.is_mac = sys.platform == 'darwin'
        self.bg_color = "#f4f6f7"
        self.root.configure(padx=10, pady=10, bg=self.bg_color) 

        style = ttk.Style()
        style.theme_use('clam')
        
        style.configure('TFrame', background=self.bg_color)
        style.configure('TLabel', background=self.bg_color, foreground="black")
        style.configure('TEntry', fieldbackground="white", foreground="black", insertcolor="black")
        style.configure("Treeview.Heading", font=("Helvetica", 10, "bold"), background="#ecf0f1", foreground="#2c3e50")
        style.configure("Treeview", rowheight=25, font=("Helvetica", 10))

        if not os.path.exists(self.db_name):
            messagebox.showwarning("Warning", "No database found! Run the generator to collect data.")
        
        self.create_widgets()
        self.refresh_all_data()

    def get_btn_style(self, hex_color):
        """Dynamically styles buttons based on Operating System"""
        if self.is_mac:
            return {"highlightbackground": hex_color, "fg": "black"}
        else:
            return {"bg": hex_color, "fg": "white"}

    def create_widgets(self):
        header_frame = ttk.Frame(self.root)
        header_frame.pack(fill="x", pady=(0, 10))
        ttk.Label(header_frame, text=f"📊 Analytics: {self.display_school_name}", font=("Helvetica", 18, "bold"), foreground="#2980b9").pack(side="left", padx=10)
        tk.Button(header_frame, text="🔄 Refresh All Data", font=("Helvetica", 10, "bold"), command=self.refresh_all_data, **self.get_btn_style("#2ecc71")).pack(side="right", padx=10, ipadx=10, ipady=3)

        self.notebook = ttk.Notebook(self.root)
        self.notebook.pack(fill="both", expand=True, padx=10, pady=10)

        self.tab_analytics = ttk.Frame(self.notebook)
        self.tab_data = ttk.Frame(self.notebook)
        self.tab_comm = ttk.Frame(self.notebook)

        self.notebook.add(self.tab_analytics, text="📊 Analytics Dashboard")
        self.notebook.add(self.tab_data, text="📋 Student Data Viewer")
        self.notebook.add(self.tab_comm, text="📨 Mass Communication")

        self.setup_analytics_tab()
        self.setup_data_tab()
        self.setup_comm_tab()

    def setup_data_tab(self):
        c_frame = ttk.Frame(self.tab_data)
        c_frame.pack(fill="x", pady=(10, 5), padx=10)

        del_frame = ttk.Frame(c_frame)
        del_frame.pack(side="left")
        
        ttk.Label(del_frame, text="🔍 Search ERP: ", font=("Helvetica", 10, "bold")).pack(side="left")
        self.delete_erp_var = tk.StringVar()
        self.delete_erp_var.trace_add("write", self.filter_treeview)
        ttk.Entry(del_frame, textvariable=self.delete_erp_var, width=15).pack(side="left", padx=(0, 10))
        tk.Button(del_frame, text="🗑️ Delete Selected Row", font=("Helvetica", 9, "bold"), command=self.delete_record, **self.get_btn_style("#e74c3c")).pack(side="left", ipadx=5)

        tk.Button(c_frame, text="💾 Export CSV", font=("Helvetica", 10, "bold"), command=self.export_csv, **self.get_btn_style("#3498db")).pack(side="right", ipadx=5, ipady=2)

        self.cols_config = [
            ("ID", 40, "center"), ("Timestamp", 130, "center"), ("Name", 130, "w"),
            ("Age", 40, "center"), ("Gender", 50, "center"), ("ERP", 80, "center"),
            ("Class", 60, "center"), ("Date", 80, "center"), ("Height", 60, "center"),
            ("Weight", 60, "center"), ("BMI", 50, "center"),
            ("Observations", 200, "w"),
            ("H-Range", 80, "center"), ("H-Status", 90, "center"),
            ("W-Range", 80, "center"), ("W-Status", 90, "center"),
            ("BMI-Range", 80, "center"), ("BMI-Status", 90, "center")
        ]
        self.columns = [c[0] for c in self.cols_config]

        tree_frame = ttk.Frame(self.tab_data)
        tree_frame.pack(fill="both", expand=True, padx=10, pady=10)

        scroll_y = ttk.Scrollbar(tree_frame, orient="vertical")
        scroll_y.pack(side="right", fill="y")
        scroll_x = ttk.Scrollbar(tree_frame, orient="horizontal")
        scroll_x.pack(side="bottom", fill="x")

        self.tree = ttk.Treeview(tree_frame, columns=self.columns, show="headings", yscrollcommand=scroll_y.set, xscrollcommand=scroll_x.set)
        scroll_y.config(command=self.tree.yview)
        scroll_x.config(command=self.tree.xview)
        
        self.tree.bind("<Double-1>", self.show_student_details)

        for col, width, align in self.cols_config:
            self.tree.column(col, width=width, anchor=align)
            self.tree.heading(col, text=col, anchor="center")

        self.tree.pack(fill="both", expand=True)

    def filter_treeview(self, *args):
        search_term = self.delete_erp_var.get().strip().lower()
        self.tree.delete(*self.tree.get_children())
        for row in self.all_student_data:
            if search_term in str(row[5]).lower(): self.tree.insert("", "end", values=row)

    def parse_dietary_data(self, raw_text):
        questions = [
            "Parent Name", "Phone Number", "Email", "Child Name", "Age", "Gender",
            "Fruits Consumed Regularly", "Vegetables Consumed Regularly",
            "Grains & Cereals Commonly Consumed", "Millets Consumed", "Milk & Dairy Intake",
            "What type of diet does your child follow?",
            "If Non-Vegetarian/Eggetarian, what does your child usually consume?",
            "Child's Eating Behaviour", "Food Preference Style", "Sugar Cravings",
            "Hydration", "Food Allergies / Intolerances", "Physical Activity",
            "Any Health Condition", "Medicine/Prescription/Supplement",
            "How did you get to know about us?"
        ]

        text = raw_text
        text = re.sub(r"[A-Z][a-z]+day, [A-Z][a-z]+ \d{1,2}, \d{4}", "", text)
        artifacts = [
            "Child Dietary Assessment Form", "Jotform", "Now create your own", 
            "PDF document", "--- PAGE", "- It's Free, Create your own PDF Document",
            "- It's Free Create your own PDF Document", "It's Free", "Create your own PDF Document", "Create your own"
        ]
        for art in artifacts: text = text.replace(art, "")

        replacements = {
            "Cer eals": "Cereals", "Dair y": "Dairy", "Intak e": "Intake",
            "y our": "your", "Non-V egetarian": "Non-Vegetarian", "E ggetarian": "Eggetarian",
            "Child' s": "Child's", "Beha viour": "Behaviour", "Pr eference": "Preference",
            "Cr avings": "Cravings", "Aller gies": "Allergies", "Int olerances": "Intolerances",
            "Pr escription": "Prescription", "y ou": "you", "t o": "to",
            "Dietar y": "Dietary", "F orm": "Form"
        }
        for bad, good in replacements.items(): text = text.replace(bad, good)

        text = re.sub(r"Grains & Cereals Commonly\s+(.*?)\s+Consumed", r"Grains & Cereals Commonly Consumed\n\1", text, flags=re.DOTALL)
        text = re.sub(r"What type of diet does your child\s+(.*?)\s+follow\??", r"What type of diet does your child follow?\n\1", text, flags=re.DOTALL)
        text = re.sub(r"If Non-Vegetarian/Eggetarian, what\s+(.*?)\s+does your child usually consume\??", r"If Non-Vegetarian/Eggetarian, what does your child usually consume?\n\1", text, flags=re.DOTALL)

        sorted_questions = sorted(questions, key=len, reverse=True)
        for q in sorted_questions:
            pattern = rf"(?<![a-zA-Z]){re.escape(q)}(?![a-zA-Z])"
            text = re.sub(pattern, f"\n{q}\n", text)

        lines = []
        for line in text.split('\n'):
            line_stripped = line.strip()
            if not line_stripped: continue
            if "Jotform" in line_stripped or "It's Free" in line_stripped or "Create your own" in line_stripped: continue
            if line_stripped.isdigit() and len(line_stripped) < 3: continue 
            line_stripped = re.sub(r'\s{2,}', ', ', line_stripped)
            lines.append(line_stripped)

        data_dict = {q: "" for q in questions}
        current_key = "Other Details"
        data_dict[current_key] = ""
            
        header_keys = ["Parent Name", "Phone Number", "Email", "Child Name", "Age", "Gender"]
        if len(lines) >= 12 and all(lines[i] == header_keys[i] for i in range(6)):
            for idx, key in enumerate(header_keys): data_dict[key] = lines[6 + idx]
            lines = lines[12:]
            
        for line in lines:
            if line in questions:
                current_key = line
            else:
                if data_dict[current_key]: data_dict[current_key] += " " + line
                else: data_dict[current_key] = line

        final_dict = {k: v.strip() for k, v in data_dict.items() if v.strip()}
        if "Other Details" in final_dict and not final_dict["Other Details"]:
            del final_dict["Other Details"]
            
        return final_dict

    def import_dietary_recall(self, erp, student_name, parent_window):
        if not HAS_FITZ: return messagebox.showerror("Missing Library", "Please install PyMuPDF to extract dietary recall forms accurately.\n\nOpen your terminal and run:\npip install PyMuPDF")

        file_path = filedialog.askopenfilename(
            title=f"Select Dietary Recall PDF for {student_name}", 
            filetypes=[("PDF Files", "*.pdf")]
        )
        if not file_path: return

        self.root.config(cursor="wait"); self.root.update()

        try:
            doc = fitz.open(file_path)
            extracted_text = ""
            for page in doc: extracted_text += page.get_text("text", sort=True) + "\n"
            doc.close()

            if not extracted_text.strip(): raise Exception("Could not extract any text. The PDF might be a scanned image.")

            clean_lines = []
            for line in extracted_text.split('\n'):
                line_stripped = line.strip()
                if not line_stripped: continue
                if "Jotform" in line_stripped or "Now create your own" in line_stripped or "It's Free" in line_stripped or "Create your own" in line_stripped or "--- PAGE" in line_stripped: continue
                if line_stripped.isdigit() and len(line_stripped) < 3: continue
                clean_lines.append(line_stripped)
            
            clean_text = "\n".join(clean_lines)

            db_filename = os.path.basename(self.db_name)
            school_name = os.path.splitext(db_filename)[0]
            safe_student_name = "".join([c for c in student_name if c.isalpha() or c.isdigit() or c==' ']).rstrip()
            safe_erp = "".join([c for c in str(erp) if c.isalpha() or c.isdigit()]).rstrip()
            
            target_dir = os.path.join(os.getcwd(), "Reports", school_name, "Diet Plans", f"{safe_student_name}-{safe_erp}")
            os.makedirs(target_dir, exist_ok=True)
            dest_file_path = os.path.join(target_dir, os.path.basename(file_path))
            shutil.copy2(file_path, dest_file_path)

            with sqlite3.connect(self.db_name) as conn:
                cursor = conn.cursor()
                cursor.execute('''CREATE TABLE IF NOT EXISTS dietary_records (
                        id INTEGER PRIMARY KEY AUTOINCREMENT, erp TEXT, import_date DATETIME DEFAULT CURRENT_TIMESTAMP, 
                        file_name TEXT, raw_extracted_text TEXT)''')
                cursor.execute('''INSERT INTO dietary_records (erp, file_name, raw_extracted_text) VALUES (?, ?, ?)''', (erp, os.path.basename(file_path), extracted_text))
                conn.commit()

            messagebox.showinfo("Success", f"Dietary Recall imported perfectly for {student_name}!")
            parent_window.destroy()
            for item in self.tree.get_children():
                if str(self.tree.item(item)['values'][5]) == str(erp):
                    self.tree.selection_set(item)
                    self.show_student_details(None)
                    break

        except Exception as e: messagebox.showerror("Import Error", f"Failed to import Dietary Recall:\n\n{str(e)}")
        finally: self.root.config(cursor="")

    # --- TRUE WYSIWYG EDITOR LOGIC & SERIALIZATION ---
    def set_active_widget(self, event):
        self.active_text_widget = event.widget

    def apply_bullet(self):
        tw = self.active_text_widget
        if not tw: return
        try:
            start_line = int(float(tw.index("sel.first")))
            end_line = int(float(tw.index("sel.last")))
            for line in range(end_line, start_line - 1, -1):
                tw.insert(f"{line}.0", "• ")
        except tk.TclError:
            current_line = int(float(tw.index("insert")))
            tw.insert(f"{current_line}.0", "• ")
        self.trigger_autosave()

    def apply_numbered_list(self):
        tw = self.active_text_widget
        if not tw: return
        try:
            start_line = int(float(tw.index("sel.first")))
            end_line = int(float(tw.index("sel.last")))
            count = 1
            for line in range(start_line, end_line + 1):
                tw.insert(f"{line}.0", f"{count}. ")
                count += 1
        except tk.TclError:
            current_line = int(float(tw.index("insert")))
            tw.insert(f"{current_line}.0", "1. ")
        self.trigger_autosave()

    def toggle_format(self, tag_name):
        tw = self.active_text_widget
        if not tw: return
        try:
            start = tw.index("sel.first")
            end = tw.index("sel.last")
            if tag_name in tw.tag_names(start):
                tw.tag_remove(tag_name, start, end)
            else:
                if tag_name.startswith("color_"):
                    for t in ["color_red", "color_green", "color_blue"]:
                        tw.tag_remove(t, start, end)
                tw.tag_add(tag_name, start, end)
        except tk.TclError:
            pass 
        self.trigger_autosave()

    def apply_font_family(self, event):
        tw = self.active_text_widget
        if not tw: return
        family = self.font_var.get()
        if family == "Font": return
        
        tag_name = f"font_{family.replace(' ', '_')}"
        tw.tag_configure(tag_name, font=(family, 10))
        try:
            start = tw.index("sel.first")
            end = tw.index("sel.last")
            for t in tw.tag_names(start):
                if t.startswith("font_"): tw.tag_remove(t, start, end)
            tw.tag_add(tag_name, start, end)
        except tk.TclError: pass
        self.font_var.set("Font")
        self.trigger_autosave()

    def apply_font_size(self, event):
        tw = self.active_text_widget
        if not tw: return
        size = self.size_var.get()
        if size == "Size": return
        
        tag_name = f"size_{size}"
        tw.tag_configure(tag_name, font=("Helvetica", int(size)))
        try:
            start = tw.index("sel.first")
            end = tw.index("sel.last")
            for t in tw.tag_names(start):
                if t.startswith("size_"): tw.tag_remove(t, start, end)
            tw.tag_add(tag_name, start, end)
        except tk.TclError: pass
        self.size_var.set("Size")
        self.trigger_autosave()

    def serialize_widget(self, tw):
        allowed_tags = ["bold_style", "italic_style", "underline_style", "color_red", "color_green", "color_blue"]
        for t in tw.tag_names():
            if t.startswith("size_") or t.startswith("font_"):
                if t not in allowed_tags:
                    allowed_tags.append(t)

        dump = tw.dump("1.0", "end-1c", text=True, tag=True)
        chunks = []
        current_tags = set()

        for key, value, index in dump:
            if key == "tagon" and value in allowed_tags:
                current_tags.add(value)
            elif key == "tagoff" and value in allowed_tags:
                current_tags.discard(value)
            elif key == "text":
                if chunks and chunks[-1]["tags"] == list(current_tags):
                    chunks[-1]["text"] += value
                else:
                    chunks.append({"text": value, "tags": list(current_tags)})

        return json.dumps(chunks)

    def deserialize_widget(self, tw, json_string):
        tw.delete("1.0", "end")
        if not json_string: return

        try:
            chunks = json.loads(json_string)
            if not isinstance(chunks, list): raise ValueError
            
            for chunk in chunks:
                start = tw.index("insert")
                tw.insert("insert", chunk.get("text", ""))
                for tag in chunk.get("tags", []):
                    if tag.startswith("size_"):
                        size = int(tag.split("_")[1])
                        tw.tag_configure(tag, font=("Helvetica", size))
                    elif tag.startswith("font_"):
                        family = tag.split("_", 1)[1].replace("_", " ")
                        tw.tag_configure(tag, font=(family, 10))
                    tw.tag_add(tag, start, "insert")
        except (ValueError, json.JSONDecodeError):
            clean_text = re.sub(r'\[/?(?:b|i|u|c:[a-zA-Z]+|s:\d+|f:[^\]]+)\]', '', json_string)
            clean_text = clean_text.replace("**", "")
            tw.insert("1.0", clean_text)

    # --- AUTO-SAVE LOGIC ---
    def trigger_autosave(self):
        if getattr(self, '_autosave_timer', None):
            self.root.after_cancel(self._autosave_timer)
        self._autosave_timer = self.root.after(1500, self.save_reviews_to_db)

    def save_reviews_to_db(self):
        if not getattr(self, 'current_erp', None): return
        
        review_data = []
        for card in getattr(self, 'review_cards', []):
            h = card['header'].get().strip()
            if h == "Subject...": h = ""
            c_json = self.serialize_widget(card['content'])
            d = card['done']
            crit = card['critical']
            
            if h or c_json != "[]":
                review_data.append({"header": h, "content_json": c_json, "done": d, "critical": crit})
        
        review_comments_json = json.dumps(review_data)
        
        field_data = {}
        for key, tw in self.diet_vars.items():
            field_data[key] = self.serialize_widget(tw)

        try:
            with sqlite3.connect(self.db_name) as conn:
                cursor = conn.cursor()
                cursor.execute("SELECT id FROM sports_diet_plans WHERE erp=?", (self.current_erp,))
                if cursor.fetchone():
                    cursor.execute('''UPDATE sports_diet_plans SET 
                        activity=?, early_morning=?, breakfast=?, lunch=?, evening_snacks=?, dinner=?,
                        pre_activity=?, during_activity=?, post_activity_20m=?, post_activity_2h=?,
                        recommendations=?, review_comments=?, updated_at=CURRENT_TIMESTAMP WHERE erp=?''', 
                        (field_data.get('activity','[]'), field_data.get('early_morning','[]'), field_data.get('breakfast','[]'),
                         field_data.get('lunch','[]'), field_data.get('evening_snacks','[]'), field_data.get('dinner','[]'),
                         field_data.get('pre_activity','[]'), field_data.get('during_activity','[]'), field_data.get('post_activity_20m','[]'),
                         field_data.get('post_activity_2h','[]'), field_data.get('recommendations','[]'), review_comments_json, self.current_erp))
                else:
                    cursor.execute('''INSERT INTO sports_diet_plans (
                        erp, activity, early_morning, breakfast, lunch, evening_snacks, dinner, pre_activity, 
                        during_activity, post_activity_20m, post_activity_2h, recommendations, review_comments) 
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)''', 
                        (self.current_erp, field_data.get('activity','[]'), field_data.get('early_morning','[]'), field_data.get('breakfast','[]'),
                         field_data.get('lunch','[]'), field_data.get('evening_snacks','[]'), field_data.get('dinner','[]'),
                         field_data.get('pre_activity','[]'), field_data.get('during_activity','[]'), field_data.get('post_activity_20m','[]'),
                         field_data.get('post_activity_2h','[]'), field_data.get('recommendations','[]'), review_comments_json))
                conn.commit()
        except Exception as e:
            print(f"Silent auto-save failed: {e}")

    # --- DYNAMIC REVIEW CARD BUILDER & LOGIC ---
    def update_card_ui(self, card):
        if card['done']:
            bg_color = "#d4edda" 
        elif card['critical']:
            bg_color = "#f48fb1" 
        else:
            bg_color = "#fff9c4" 
            
        card['frame'].config(bg=bg_color)
        card['top_row'].config(bg=bg_color)
        card['header_entry'].config(bg=bg_color)
        card['content'].config(bg=bg_color)
        
        card['btn_done'].config(text="DONE" if card['done'] else "PENDING", bg="#27ae60" if card['done'] else "#95a5a6")
        card['btn_crit'].config(text="CRITICAL" if card['critical'] else "NORMAL", bg="#d81b60" if card['critical'] else "#95a5a6")

    def sort_and_repack_cards(self):
        def get_priority(c):
            if c['done']: return 2 if c['critical'] else 3
            if c['critical']: return 0
            return 1
            
        self.review_cards.sort(key=get_priority)
        for c in self.review_cards:
            c['frame'].pack_forget()
            c['frame'].pack(fill="x", pady=5, padx=5)

    def create_review_card(self, parent_frame, header_text="", content_json="", is_done=False, is_critical=False):
        card = {}
        card['done'] = is_done
        card['critical'] = is_critical
        
        card_frame = tk.Frame(parent_frame, bd=1, relief="solid")
        
        top_row = tk.Frame(card_frame)
        top_row.pack(fill="x", padx=5, pady=5)
        
        header_var = tk.StringVar(value=header_text)
        header_entry = tk.Entry(top_row, textvariable=header_var, font=("Helvetica", 10, "bold"), relief="flat", width=12)
        header_entry.pack(side="left", fill="x", expand=True)
        
        if not header_text: header_entry.insert(0, "Subject...")
        
        header_entry.bind("<FocusIn>", lambda e: header_entry.delete(0, 'end') if header_entry.get() == "Subject..." else None)
        header_entry.bind("<FocusOut>", lambda e: [header_entry.insert(0, "Subject...") if not header_entry.get() else None, self.trigger_autosave()])
        header_entry.bind("<KeyRelease>", lambda e: self.trigger_autosave())

        def toggle_done():
            card['done'] = not card['done']
            self.update_card_ui(card)
            self.sort_and_repack_cards()
            self.trigger_autosave()

        def toggle_crit():
            card['critical'] = not card['critical']
            self.update_card_ui(card)
            self.sort_and_repack_cards()
            self.trigger_autosave()

        def delete_card():
            card_frame.destroy()
            if card in self.review_cards:
                self.review_cards.remove(card)
            self.trigger_autosave()

        btn_crit = tk.Button(top_row, font=("Helvetica", 8, "bold"), command=toggle_crit, **self.get_btn_style("#95a5a6"))
        btn_crit.pack(side="left", padx=2)

        btn_done = tk.Button(top_row, font=("Helvetica", 8, "bold"), command=toggle_done, **self.get_btn_style("#95a5a6"))
        btn_done.pack(side="left", padx=2)
        
        tk.Button(top_row, text="X", font=("Helvetica", 8, "bold"), command=delete_card, **self.get_btn_style("#e74c3c")).pack(side="left", padx=(2,0))
        
        content_txt = tk.Text(card_frame, height=4, width=25, font=("Helvetica", 10), wrap="word", relief="flat")
        content_txt.pack(fill="both", expand=True, padx=5, pady=(0, 5))
        
        content_txt.tag_configure("bold_style", font=("Helvetica", 10, "bold"))
        content_txt.tag_configure("italic_style", font=("Helvetica", 10, "italic"))
        content_txt.tag_configure("underline_style", underline=True)
        content_txt.tag_configure("color_red", foreground="red")
        content_txt.tag_configure("color_green", foreground="green")
        content_txt.tag_configure("color_blue", foreground="blue")
        
        self.deserialize_widget(content_txt, content_json)
        
        content_txt.bind("<FocusIn>", self.set_active_widget)
        content_txt.bind("<KeyRelease>", lambda e: self.trigger_autosave())
        
        card['frame'] = card_frame
        card['top_row'] = top_row
        card['header_entry'] = header_entry
        card['content'] = content_txt
        card['header'] = header_var
        card['btn_done'] = btn_done
        card['btn_crit'] = btn_crit
        
        self.review_cards.append(card)
        self.update_card_ui(card)
        self.sort_and_repack_cards()
        parent_frame.update_idletasks()
        
        if hasattr(self, 'current_erp') and self.current_erp:
            self.trigger_autosave()

    def setup_diet_planning_tab(self, parent_frame, erp, name, age, gender):
        for widget in parent_frame.winfo_children():
            widget.destroy()

        btn_frame = tk.Frame(parent_frame, bg="#f4f6f7")
        btn_frame.pack(expand=True, fill="both")
        
        inner_frame = tk.Frame(btn_frame, bg="#f4f6f7")
        inner_frame.place(relx=0.5, rely=0.5, anchor="center")

        has_sports_plan = False
        sports_updated = ""
        try:
            with sqlite3.connect(self.db_name) as conn:
                cursor = conn.cursor()
                cursor.execute("SELECT updated_at FROM sports_diet_plans WHERE erp=?", (erp,))
                row = cursor.fetchone()
                if row:
                    has_sports_plan = True
                    sports_updated = str(row[0]).split('.')[0]
        except sqlite3.OperationalError: pass

        tk.Button(inner_frame, text="Create Regular Diet Plan", font=("Helvetica", 13, "bold"), state="disabled", width=40, **self.get_btn_style("#bdc3c7")).pack(pady=15, ipady=8)
        
        sports_btn_text = f"✏️ Edit Sports Plan (Saved: {sports_updated})" if has_sports_plan else "Create Sports Diet Plan"
        sports_btn_bg = "#27ae60" if has_sports_plan else "#e67e22"

        tk.Button(inner_frame, text=sports_btn_text, font=("Helvetica", 13, "bold"), width=40, command=lambda: self.show_sports_plan_form(parent_frame, erp, name, age, gender), **self.get_btn_style(sports_btn_bg)).pack(pady=15, ipady=8)

    def show_sports_plan_form(self, parent_frame, erp, name, age, gender):
        for widget in parent_frame.winfo_children():
            widget.destroy()
            
        self.review_cards = [] 
        self.current_erp = erp 

        # --- TRUE WYSIWYG TOOLBAR ---
        toolbar_frame = tk.Frame(parent_frame, bg="#ecf0f1", bd=1, relief="solid")
        toolbar_frame.pack(side="top", fill="x")
        
        tk.Button(toolbar_frame, text="← Back", font=("Helvetica", 9, "bold"), command=lambda: self.setup_diet_planning_tab(parent_frame, erp, name, age, gender), **self.get_btn_style("#95a5a6")).pack(side="left", padx=10, pady=5)
        ttk.Separator(toolbar_frame, orient='vertical').pack(side="left", fill="y", padx=5)
        
        # Core Styles
        tk.Button(toolbar_frame, text=" B ", font=("Helvetica", 10, "bold"), relief="groove", command=lambda: self.toggle_format("bold_style")).pack(side="left", padx=2, pady=5)
        tk.Button(toolbar_frame, text=" I ", font=("Helvetica", 10, "italic"), relief="groove", command=lambda: self.toggle_format("italic_style")).pack(side="left", padx=2, pady=5)
        tk.Button(toolbar_frame, text=" U ", font=("Helvetica", 10, "underline"), relief="groove", command=lambda: self.toggle_format("underline_style")).pack(side="left", padx=2, pady=5)
        ttk.Separator(toolbar_frame, orient='vertical').pack(side="left", fill="y", padx=5)

        # Colors
        tk.Button(toolbar_frame, text=" Red ", fg="red", relief="groove", command=lambda: self.toggle_format("color_red")).pack(side="left", padx=2, pady=5)
        tk.Button(toolbar_frame, text=" Grn ", fg="green", relief="groove", command=lambda: self.toggle_format("color_green")).pack(side="left", padx=2, pady=5)
        tk.Button(toolbar_frame, text=" Blu ", fg="blue", relief="groove", command=lambda: self.toggle_format("color_blue")).pack(side="left", padx=2, pady=5)
        ttk.Separator(toolbar_frame, orient='vertical').pack(side="left", fill="y", padx=5)

        # Lists
        tk.Button(toolbar_frame, text=" • ", font=("Helvetica", 10, "bold"), relief="groove", command=self.apply_bullet).pack(side="left", padx=2, pady=5)
        tk.Button(toolbar_frame, text=" 1. ", font=("Helvetica", 10, "bold"), relief="groove", command=self.apply_numbered_list).pack(side="left", padx=2, pady=5)
        ttk.Separator(toolbar_frame, orient='vertical').pack(side="left", fill="y", padx=5)

        # Dynamic Fonts
        self.font_var = tk.StringVar(value="Font")
        font_cb = ttk.Combobox(toolbar_frame, textvariable=self.font_var, values=["Arial", "Calibri", "Times New Roman"], width=15, state="readonly")
        font_cb.pack(side="left", padx=5, pady=5)
        font_cb.bind("<<ComboboxSelected>>", self.apply_font_family)

        self.size_var = tk.StringVar(value="Size")
        size_cb = ttk.Combobox(toolbar_frame, textvariable=self.size_var, values=["10", "12", "14", "18", "24"], width=5, state="readonly")
        size_cb.pack(side="left", padx=5, pady=5)
        size_cb.bind("<<ComboboxSelected>>", self.apply_font_size)

        try:
            with sqlite3.connect(self.db_name) as conn:
                cursor = conn.cursor()
                cursor.execute('''CREATE TABLE IF NOT EXISTS sports_diet_plans (
                    id INTEGER PRIMARY KEY AUTOINCREMENT, erp TEXT UNIQUE,
                    updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                    activity TEXT, early_morning TEXT, breakfast TEXT,
                    lunch TEXT, evening_snacks TEXT, dinner TEXT,
                    pre_activity TEXT, during_activity TEXT,
                    post_activity_20m TEXT, post_activity_2h TEXT,
                    recommendations TEXT)''')
                try: cursor.execute("ALTER TABLE sports_diet_plans ADD COLUMN review_comments TEXT")
                except sqlite3.OperationalError: pass
        except Exception: pass

        existing_plan = {}
        try:
            with sqlite3.connect(self.db_name) as conn:
                cursor = conn.cursor()
                cursor.execute("SELECT * FROM sports_diet_plans WHERE erp=?", (erp,))
                row = cursor.fetchone()
                if row:
                    cols = [col[0] for col in cursor.description]
                    existing_plan = dict(zip(cols, row))
        except sqlite3.OperationalError: pass

        self.diet_vars = {}

        content_frame = tk.Frame(parent_frame, bg="#f4f6f7")
        content_frame.pack(fill="both", expand=True)

        sticky_container = tk.Frame(content_frame, bg="#fef9e7", bd=2, relief="groove", width=340)
        sticky_container.pack(side="right", fill="y", padx=15, pady=15)
        sticky_container.pack_propagate(False) 
        
        sticky_header = tk.Frame(sticky_container, bg="#fef9e7")
        sticky_header.pack(fill="x", padx=10, pady=(10, 5))
        
        tk.Label(sticky_header, text="📌 Review Comments", font=("Helvetica", 12, "bold"), bg="#fef9e7", fg="#d35400").pack(side="left")
        tk.Label(sticky_container, text="(Internal use only)", font=("Helvetica", 9, "italic"), bg="#fef9e7", fg="gray").pack(anchor="w", padx=10, pady=(0, 5))
        
        sc_canvas = tk.Canvas(sticky_container, bg="#fef9e7", highlightthickness=0)
        sc_scrollbar = ttk.Scrollbar(sticky_container, orient="vertical", command=sc_canvas.yview)
        sticky_scroll_frame = tk.Frame(sc_canvas, bg="#fef9e7")
        
        sticky_scroll_frame.bind("<Configure>", lambda e: sc_canvas.configure(scrollregion=sc_canvas.bbox("all")))
        sc_canvas.create_window((0, 0), window=sticky_scroll_frame, anchor="nw", width=300)
        sc_canvas.configure(yscrollcommand=sc_scrollbar.set)
        
        sc_canvas.pack(side="left", fill="both", expand=True, pady=5, padx=5)
        sc_scrollbar.pack(side="right", fill="y", pady=5)
        
        tk.Button(sticky_header, text=" + ", font=("Helvetica", 10, "bold"), command=lambda: self.create_review_card(sticky_scroll_frame), **self.get_btn_style("#2980b9")).pack(side="right")

        initial_load_erp = self.current_erp
        self.current_erp = None 
        
        if existing_plan and existing_plan.get("review_comments"):
            try:
                saved_cards = json.loads(existing_plan["review_comments"])
                if isinstance(saved_cards, list) and len(saved_cards) > 0:
                    for sc in saved_cards:
                        self.create_review_card(sticky_scroll_frame, sc.get("header", ""), sc.get("content_json", sc.get("content", "")), sc.get("done", False), sc.get("critical", False))
                else:
                    self.create_review_card(sticky_scroll_frame, "Legacy Note", str(existing_plan["review_comments"]), False, False)
            except Exception:
                self.create_review_card(sticky_scroll_frame, "Legacy Note", str(existing_plan["review_comments"]), False, False)
        else:
            self.create_review_card(sticky_scroll_frame, "Initial Review", "", False, False) 
            
        self.current_erp = initial_load_erp

        canvas = tk.Canvas(content_frame, bg="#f4f6f7", highlightthickness=0)
        scrollbar = ttk.Scrollbar(content_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg="#f4f6f7", padx=20, pady=20)

        scrollable_frame.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="left", fill="y")

        ttk.Label(scrollable_frame, text=f"🏋️ Sports Diet Plan Builder: {name}", font=("Helvetica", 16, "bold"), background="#f4f6f7", foreground="#2c3e50").grid(row=0, column=0, columnspan=2, sticky="w", pady=(0, 20))

        fields = [
            ("Activity / Sport", "activity", 1),
            ("Early Morning", "early_morning", 3),
            ("Breakfast / Post-Workout", "breakfast", 4),
            ("Lunch", "lunch", 4),
            ("Evening Snacks", "evening_snacks", 4),
            ("Dinner", "dinner", 4),
            ("Pre-Activity Meal (30 mins prior)", "pre_activity", 3),
            ("During Activity (Hydration/Snack)", "during_activity", 3),
            ("Post-Activity (Within 20 mins)", "post_activity_20m", 3),
            ("Post-Activity (Within 2 hours)", "post_activity_2h", 3),
            ("Dietary Recommendations", "recommendations", 8)
        ]

        row_idx = 1
        for label_text, key, height in fields:
            ttk.Label(scrollable_frame, text=label_text + ":", font=("Helvetica", 10, "bold"), foreground="#2c3e50", background="#f4f6f7").grid(row=row_idx, column=0, sticky="nw", pady=8)
            
            txt_widget = tk.Text(scrollable_frame, height=height, width=65, font=("Helvetica", 10), wrap="word", relief="solid", borderwidth=1, bg="white", fg="black", insertbackground="black")
            txt_widget.grid(row=row_idx, column=1, sticky="w", pady=8, padx=15)
            
            txt_widget.tag_configure("bold_style", font=("Helvetica", 10, "bold"))
            txt_widget.tag_configure("italic_style", font=("Helvetica", 10, "italic"))
            txt_widget.tag_configure("underline_style", underline=True)
            txt_widget.tag_configure("color_red", foreground="red")
            txt_widget.tag_configure("color_green", foreground="green")
            txt_widget.tag_configure("color_blue", foreground="blue")
            
            txt_widget.bind("<FocusIn>", self.set_active_widget)
            txt_widget.bind("<KeyRelease>", lambda e: self.trigger_autosave())
            
            if key in existing_plan and existing_plan[key]:
                self.deserialize_widget(txt_widget, existing_plan[key])
                
            self.diet_vars[key] = txt_widget
            row_idx += 1

        btn_text = "Update & Export Sports Plan (PDF)" if existing_plan else "Save & Export Sports Plan (PDF)"
        tk.Button(scrollable_frame, text=btn_text, font=("Helvetica", 12, "bold"), command=lambda: self.generate_sports_plan(erp, name, age, gender), **self.get_btn_style("#2ecc71")).grid(row=row_idx, column=0, columnspan=2, pady=30, ipadx=20, ipady=8)

    def generate_sports_plan(self, erp, student_name, age, gender):
        if not HAS_PPTX:
            return messagebox.showerror("Missing Library", "The 'python-pptx' library is required to generate Diet Plans.\n\nOpen your terminal and run:\npip install python-pptx")
            
        template_path = os.path.join(os.getcwd(), "template", "sports_diet_template.pptx")
        if not os.path.exists(template_path):
            return messagebox.showerror("Missing Template", f"Could not find the template file at:\n{template_path}\n\nPlease place your PPTX file there and name it 'sports_diet_template.pptx'.")
            
        self.root.config(cursor="wait"); self.root.update()

        db_filename = os.path.basename(self.db_name)
        school_name = os.path.splitext(db_filename)[0]
        
        first_name = student_name.strip().split()[0]
        safe_student_name = "".join([c for c in student_name if c.isalpha() or c.isdigit() or c==' ']).rstrip()
        safe_erp = "".join([c for c in str(erp) if c.isalpha() or c.isdigit()]).rstrip()
        
        target_dir = os.path.join(os.getcwd(), "Reports", school_name, "Diet Plans", f"{safe_student_name}-{safe_erp}")
        os.makedirs(target_dir, exist_ok=True)
        
        temp_pptx = os.path.join(target_dir, f"temp_{safe_student_name}_Sports_Plan.pptx")
        output_pdf = os.path.join(target_dir, f"{first_name}_Sports_Plan.pdf")
        
        self.trigger_autosave() 
        
        try:
            prs = pptx.Presentation(template_path)
            
            basic_replacements = {
                "{{name}}": first_name,
                "{{age}}": str(age),
                "{{gender}}": "Female" if gender == "F" else "Male"
            }
            
            ast_replacements = {}
            for key, txt_widget in self.diet_vars.items():
                tag = "{{" + key + "}}"
                ast_replacements[tag] = self.serialize_widget(txt_widget)

            def apply_json_to_paragraph(paragraph, json_chunks, base_font_name, base_font_size, base_font_bold, base_font_italic, base_font_color):
                for chunk in json_chunks:
                    run = paragraph.add_run()
                    run.text = chunk.get("text", "")
                    tags = chunk.get("tags", [])

                    if base_font_name: run.font.name = base_font_name
                    if base_font_size: run.font.size = base_font_size
                    run.font.bold = base_font_bold
                    run.font.italic = base_font_italic
                    if base_font_color:
                        try: run.font.color.rgb = base_font_color
                        except: pass

                    if "bold_style" in tags: run.font.bold = True
                    if "italic_style" in tags: run.font.italic = True
                    if "underline_style" in tags: run.font.underline = True
                    if "color_red" in tags: run.font.color.rgb = pptx.dml.color.RGBColor(255,0,0)
                    if "color_green" in tags: run.font.color.rgb = pptx.dml.color.RGBColor(0,128,0)
                    if "color_blue" in tags: run.font.color.rgb = pptx.dml.color.RGBColor(0,0,255)

                    for t in tags:
                        if t.startswith("size_"):
                            run.font.size = pptx.util.Pt(int(t.split("_")[1]))
                        if t.startswith("font_"):
                            run.font.name = t.split("_", 1)[1].replace("_", " ")

            def replace_in_shape(shape):
                if getattr(shape, "has_text_frame", False) or hasattr(shape, "text_frame"):
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for tag, value in basic_replacements.items():
                                if tag in paragraph.text:
                                    paragraph.text = paragraph.text.replace(tag, value)
                            
                            p_text = paragraph.text
                            for tag, json_val in ast_replacements.items():
                                if tag in p_text:
                                    font_name, font_size, font_bold, font_italic, font_color = None, None, False, False, None
                                    if paragraph.runs:
                                        r = paragraph.runs[0]
                                        font_name, font_size = r.font.name, r.font.size
                                        font_bold, font_italic = r.font.bold, r.font.italic
                                        if hasattr(r.font.color, 'type') and r.font.color.type is not None:
                                            try: font_color = r.font.color.rgb
                                            except: pass
                                    
                                    try:
                                        chunks = json.loads(json_val)
                                        if p_text.strip() == tag:
                                            paragraph.text = "" 
                                            apply_json_to_paragraph(paragraph, chunks, font_name, font_size, font_bold, font_italic, font_color)
                                        else:
                                            parts = p_text.split(tag)
                                            paragraph.text = ""
                                            if parts[0]: paragraph.add_run().text = parts[0]
                                            apply_json_to_paragraph(paragraph, chunks, font_name, font_size, font_bold, font_italic, font_color)
                                            if len(parts) > 1 and parts[1]: paragraph.add_run().text = parts[1]
                                            p_text = paragraph.text 
                                    except json.JSONDecodeError:
                                        paragraph.text = paragraph.text.replace(tag, json_val)

                if getattr(shape, "shape_type", None) == 6:
                    for child in shape.shapes:
                        replace_in_shape(child)
                
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            replace_in_shape(cell)

            for slide in prs.slides:
                for shape in slide.shapes:
                    replace_in_shape(shape)
                                
            prs.save(temp_pptx)
            
            pdf_msg = ""
            if os.name == 'nt' and HAS_COMTYPES:
                try:
                    import comtypes.client
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    deck = powerpoint.Presentations.Open(os.path.abspath(temp_pptx), WithWindow=False)
                    deck.SaveAs(os.path.abspath(output_pdf), 32)
                    deck.Close()
                    powerpoint.Quit()
                    pdf_msg = "\n\n✓ PDF successfully generated!"
                except Exception as e:
                    pdf_msg = f"\n\n⚠ Temporary PPTX saved, but PDF conversion failed: {e}"
            elif sys.platform == 'darwin':
                try:
                    import subprocess
                    script = f'''
                    tell application "Microsoft PowerPoint"
                        open POSIX file "{os.path.abspath(temp_pptx)}"
                        save active presentation in POSIX file "{os.path.abspath(output_pdf)}" as save as PDF
                        close active presentation saving no
                    end tell
                    '''
                    subprocess.run(['osascript', '-e', script], check=True)
                    pdf_msg = "\n\n✓ PDF successfully generated using Mac PowerPoint!"
                except Exception as e:
                    pdf_msg = f"\n\n⚠ PPTX saved, but Mac PDF conversion failed (Requires MS PowerPoint installed natively): {e}"
            else:
                pdf_msg = "\n\n⚠ Automatic PDF conversion requires Windows OS or macOS with MS Office. A PPTX was saved instead."
                
            if not (os.name == 'nt' and HAS_COMTYPES) and sys.platform != 'darwin':
                final_pptx = os.path.join(target_dir, f"{first_name}_Sports_Plan.pptx")
                shutil.move(temp_pptx, final_pptx)
                temp_pptx = None
                
            if temp_pptx and os.path.exists(temp_pptx):
                try: os.remove(temp_pptx)
                except Exception as e: print(f"Warning: Could not remove temporary PPTX file: {e}")
                
            messagebox.showinfo("Success", f"Sports Diet Plan generated and exported successfully!\n\nLocation:\n{target_dir}{pdf_msg}")
            
        except Exception as e:
            messagebox.showerror("Generation Error", f"An error occurred while generating the plan:\n{str(e)}")
        finally:
            self.root.config(cursor="")

    def show_student_details(self, event):
        selected = self.tree.selection()
        if not selected: return
        item = self.tree.item(selected[0])
        values = item['values']
        
        name, age, gender, erp, cls = values[2], values[3], values[4], values[5], values[6]
        height, weight, bmi, obs = values[8], values[9], values[10], values[11]
        h_stat, w_stat, bmi_stat = values[13], values[15], values[17]

        top = tk.Toplevel(self.root)
        top.title(f"Student Health Profile - {name}")
        top.geometry("1000x700") 
        top.configure(bg="#f4f6f7", padx=20, pady=20)
        
        ttk.Label(top, text=f"🎓 {name}", font=("Helvetica", 20, "bold"), background="#f4f6f7", foreground="#2c3e50").pack(anchor="w")
        ttk.Label(top, text=f"ERP: {erp}  |  Class: {cls}  |  Age: {age}  |  Gender: {gender}", font=("Helvetica", 11, "bold"), background="#f4f6f7", foreground="#7f8c8d").pack(anchor="w", pady=(5, 10))
        
        profile_notebook = ttk.Notebook(top)
        profile_notebook.pack(fill="both", expand=True, pady=10)

        tab_health = ttk.Frame(profile_notebook)
        tab_diet = ttk.Frame(profile_notebook)
        tab_plan = ttk.Frame(profile_notebook)

        profile_notebook.add(tab_health, text="🩺 Health & Growth Profile")
        profile_notebook.add(tab_diet, text="🥗 Dietary Recall Data")
        profile_notebook.add(tab_plan, text="📝 Diet Planning")

        self.setup_diet_planning_tab(tab_plan, erp, name, age, gender)

        # ==========================================
        # TAB 1: HEALTH PROFILE
        # ==========================================
        metrics_frame = ttk.LabelFrame(tab_health, text=" Recorded Health Metrics ", padding=15)
        metrics_frame.pack(fill="x", padx=15, pady=(15, 10))
        
        def add_metric(parent, row, label, val, stat):
            ttk.Label(parent, text=label, font=("Helvetica", 11, "bold")).grid(row=row, column=0, sticky="w", pady=5, padx=5)
            ttk.Label(parent, text=str(val), font=("Helvetica", 11)).grid(row=row, column=1, sticky="w", pady=5, padx=20)
            color = "#27ae60" if str(stat).lower() == "normal" else "#e67e22"
            stat_lbl = tk.Label(parent, text=str(stat).upper(), font=("Helvetica", 9, "bold"), bg=color, fg="white", padx=8, pady=3, borderwidth=1, relief="solid")
            stat_lbl.grid(row=row, column=2, sticky="w", pady=5, padx=5)

        add_metric(metrics_frame, 0, "Height (cm):", height, h_stat)
        add_metric(metrics_frame, 1, "Weight (kg):", weight, w_stat)
        add_metric(metrics_frame, 2, "Calculated BMI:", bmi, bmi_stat)
        
        rda_frame = ttk.LabelFrame(tab_health, text=" Recommended Dietary Allowance (RDA) Estimates ", padding=15)
        rda_frame.pack(fill="x", padx=15, pady=(0, 10))
        
        try: a = int(str(age).split('/')[0])
        except: a = 10
        
        if a <= 3:
            rda = {"Energy": "~1110 Kcal", "Protein": "38 g", "Cereals/Millets": "100 g", "^Pulses & Beans": "50 g", "Green Leafy Veg": "50 g", "Other Veg": "100 g", "Roots & Tubers": "50 g", "Fruits": "60-75 g", "Nuts": "10 g", "Milk/Curd": "350 ml", "Fats & Oils": "20 g"}
        elif 4 <= a <= 6:
            rda = {"Energy": "~1370 Kcal", "Protein": "46 g", "Cereals/Millets": "160 g", "^Pulses & Beans": "60 g", "Green Leafy Veg": "50 g", "Other Veg": "100 g", "Roots & Tubers": "50 g", "Fruits": "75 g", "Nuts": "15 g", "Milk/Curd": "350 ml", "Fats & Oils": "20 g"}
        elif 7 <= a <= 9:
            rda = {"Energy": "~1710 Kcal", "Protein": "59 g", "Cereals/Millets": "200 g", "^Pulses & Beans": "65 g", "Green Leafy Veg": "100 g", "Other Veg": "150 g", "Roots & Tubers": "100 g", "Fruits": "100 g", "Nuts": "20 g", "Milk/Curd": "400 ml", "Fats & Oils": "25 g"}
        elif 10 <= a <= 12:
            if gender == 'M': rda = {"Energy": "~2230 Kcal", "Protein": "76 g", "Cereals/Millets": "280 g", "^Pulses & Beans": "90 g", "Green Leafy Veg": "100 g", "Other Veg": "200 g", "Roots & Tubers": "100 g", "Fruits": "100 g", "Nuts": "30 g", "Milk/Curd": "400 ml", "Fats & Oils": "35 g"}
            else: rda = {"Energy": "~2060 Kcal", "Protein": "70 g", "Cereals/Millets": "250 g", "^Pulses & Beans": "85 g", "Green Leafy Veg": "100 g", "Other Veg": "200 g", "Roots & Tubers": "100 g", "Fruits": "100 g", "Nuts": "30 g", "Milk/Curd": "400 ml", "Fats & Oils": "30 g"}
        elif 13 <= a <= 15:
            if gender == 'M': rda = {"Energy": "~2860 Kcal", "Protein": "95 g", "Cereals/Millets": "390 g", "^Pulses & Beans": "130 g", "Green Leafy Veg": "100 g", "Other Veg": "200 g", "Roots & Tubers": "100 g", "Fruits": "100 g", "Nuts": "40 g", "Milk/Curd": "400 ml", "Fats & Oils": "45 g"}
            else: rda = {"Energy": "~2410 Kcal", "Protein": "81 g", "Cereals/Millets": "300 g", "^Pulses & Beans": "100 g", "Green Leafy Veg": "100 g", "Other Veg": "200 g", "Roots & Tubers": "100 g", "Fruits": "100 g", "Nuts": "35 g", "Milk/Curd": "400 ml", "Fats & Oils": "40 g"}
        else:
            if gender == 'M': rda = {"Energy": "~3300 Kcal", "Protein": "107 g", "Cereals/Millets": "450 g", "^Pulses & Beans": "150 g", "Green Leafy Veg": "100 g", "Other Veg": "200 g", "Roots & Tubers": "100 g", "Fruits": "150 g", "Nuts": "50 g", "Milk/Curd": "400 ml", "Fats & Oils": "55 g"}
            else: rda = {"Energy": "~2490 Kcal", "Protein": "85 g", "Cereals/Millets": "315 g", "^Pulses & Beans": "105 g", "Green Leafy Veg": "100 g", "Other Veg": "200 g", "Roots & Tubers": "100 g", "Fruits": "150 g", "Nuts": "40 g", "Milk/Curd": "400 ml", "Fats & Oils": "40 g"}

        row_idx = 0
        col_idx = 0
        for k, v in rda.items():
            ttk.Label(rda_frame, text=f"• {k}:", font=("Helvetica", 10, "bold")).grid(row=row_idx, column=col_idx*2, sticky="w", pady=4, padx=5)
            ttk.Label(rda_frame, text=v, font=("Helvetica", 10)).grid(row=row_idx, column=col_idx*2+1, sticky="w", pady=4, padx=(0, 25))
            row_idx += 1
            if row_idx >= 6: 
                row_idx = 0
                col_idx += 1

        ttk.Label(rda_frame, text="^ For non-vegetarians, 30g of pulses may be substituted with meat or eggs.", font=("Helvetica", 9, "italic"), foreground="#7f8c8d").grid(row=6, column=0, columnspan=4, sticky="w", pady=(10, 0), padx=5)

        obs_frame = ttk.LabelFrame(tab_health, text=" Generated Observations ", padding=15)
        obs_frame.pack(fill="both", expand=True, padx=15, pady=(0, 15))
        
        txt = tk.Text(obs_frame, font=("Helvetica", 10), wrap="word", bg="#ecf0f1", relief="flat", padx=10, pady=10)
        txt.pack(fill="both", expand=True)
        txt.insert("1.0", str(obs))
        txt.config(state="disabled")

        # ==========================================
        # TAB 2: DIETARY RECALL
        # ==========================================
        diet_action_frame = ttk.Frame(tab_diet)
        diet_action_frame.pack(fill="x", pady=15, padx=15)
        
        ttk.Label(diet_action_frame, text="Upload and extract data from a Jotform Dietary Recall PDF.", foreground="grey").pack(side="left")
        tk.Button(diet_action_frame, text="📄 Import Dietary PDF", font=("Helvetica", 10, "bold"), command=lambda: self.import_dietary_recall(erp, name, top), **self.get_btn_style("#8e44ad")).pack(side="right", ipadx=10, ipady=3)
        
        try:
            with sqlite3.connect(self.db_name) as conn:
                cursor = conn.cursor()
                cursor.execute("SELECT import_date, file_name, raw_extracted_text FROM dietary_records WHERE erp = ? ORDER BY id DESC", (erp,))
                diet_records = cursor.fetchall()
        except sqlite3.OperationalError:
            diet_records = []
            
        diet_content_frame = ttk.LabelFrame(tab_diet, text=" Extracted Dietary Data ", padding=15)
        diet_content_frame.pack(fill="both", expand=True, padx=15, pady=(0, 15))
        
        if diet_records:
            latest_date, latest_file, raw_text = diet_records[0]
            
            info_frame = tk.Frame(diet_content_frame, bg="#f4f6f7")
            info_frame.pack(fill="x", pady=(0, 10))
            ttk.Label(info_frame, text=f"File Name: {latest_file}", font=("Helvetica", 10, "bold"), background="#f4f6f7").pack(anchor="w")
            ttk.Label(info_frame, text=f"Imported On: {latest_date}", font=("Helvetica", 9, "italic"), foreground="grey", background="#f4f6f7").pack(anchor="w")
            
            diet_scroll_y = ttk.Scrollbar(diet_content_frame, orient="vertical")
            diet_scroll_y.pack(side="right", fill="y")

            diet_txt = tk.Text(diet_content_frame, font=("Helvetica", 11), wrap="word", bg="#ffffff", relief="flat", padx=10, pady=10, yscrollcommand=diet_scroll_y.set)
            diet_txt.pack(fill="both", expand=True)
            diet_scroll_y.config(command=diet_txt.yview)

            diet_txt.tag_configure("question", font=("Helvetica", 11, "bold"), foreground="#2980b9", spacing1=12)
            diet_txt.tag_configure("answer", font=("Helvetica", 11), foreground="#2c3e50", spacing3=5)

            parsed_data = self.parse_dietary_data(raw_text)
            
            if len(parsed_data) > 3:
                for q, a in parsed_data.items():
                    if not a: a = "N/A"
                    diet_txt.insert("end", f"{q}\n", "question")
                    diet_txt.insert("end", f"  • {a}\n", "answer")
            else:
                diet_txt.insert("1.0", str(raw_text))

            diet_txt.config(state="disabled")
        else:
            ttk.Label(diet_content_frame, text="No dietary recall data found for this student.", font=("Helvetica", 11, "italic"), foreground="grey").pack(pady=40)

    def export_csv(self):
        if not self.tree.get_children(): return messagebox.showwarning("No Data", "No data to export!")
        
        db_filename = os.path.basename(self.db_name)
        school_name = os.path.splitext(db_filename)[0]
        export_dir = os.path.join(os.getcwd(), "Reports", school_name, "exports", "student_data")
        os.makedirs(export_dir, exist_ok=True)
        
        safe_school_name = self.display_school_name.replace(" ", "_")
        path = os.path.join(export_dir, f"{safe_school_name}_export_data.csv")
        try:
            with open(path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                writer.writerow(self.columns)
                for item in self.tree.get_children(): writer.writerow(self.tree.item(item)['values'])
            messagebox.showinfo("Success", f"Data auto-saved successfully!\n\nLocation:\n{path}")
        except Exception as e: messagebox.showerror("Export Error", f"Failed to export data:\n{str(e)}")

    def delete_record(self):
        selected = self.tree.selection()
        if not selected: return messagebox.showwarning("Selection Needed", "Please select a row from the table below to delete.")
        item = self.tree.item(selected[0])
        student_id, student_name, erp = item['values'][0], item['values'][2], str(item['values'][5])

        if not messagebox.askyesno("Confirm Deletion", f"Found record for: {student_name} (ERP: {erp})\n\nAre you sure you want to permanently delete this student's data, their generated Word document, and their Master Report?"): return

        try:
            with sqlite3.connect(self.db_name) as conn:
                cursor = conn.cursor()
                cursor.execute("SELECT report_date FROM growth_reports WHERE id = ?", (student_id,))
                record = cursor.fetchone()
                
                if record:
                    date_folder = str(record[0]).replace('/', '-') if record[0] else 'Unknown_Date'
                    base_reports_dir = os.path.join(os.getcwd(), "Reports")
                    db_filename = os.path.basename(self.db_name)
                    school_name = os.path.splitext(db_filename)[0]

                    doc_path = os.path.join(base_reports_dir, school_name, "interpretation_docs", date_folder, f"LVLX_Growth_Report_{erp}.docx")
                    pdf_path = os.path.join(base_reports_dir, school_name, "Master Report", date_folder, f"LVLX_Master_Report_{erp}.pdf")

                    if os.path.exists(doc_path): os.remove(doc_path)
                    if os.path.exists(pdf_path): os.remove(pdf_path)

                cursor.execute("DELETE FROM growth_reports WHERE id = ?", (student_id,))
                conn.commit()

            messagebox.showinfo("Success", f"Record, document, and Master Report for {student_name} have been safely deleted.")
            self.delete_erp_var.set("")
            self.refresh_all_data()

        except Exception as e: messagebox.showerror("Database Error", str(e))

    def setup_comm_tab(self):
        self.comm_notebook = ttk.Notebook(self.tab_comm)
        self.comm_notebook.pack(fill="x", padx=20, pady=(15, 10))

        self.frame_email = ttk.Frame(self.comm_notebook, padding=15)
        self.comm_notebook.add(self.frame_email, text="📧 Email Engine Setup")
        
        ttk.Label(self.frame_email, text="Sender Email (Gmail):", font=("Helvetica", 10, "bold")).grid(row=0, column=0, sticky="w", pady=5, padx=5)
        self.sender_email_var = tk.StringVar()
        ttk.Entry(self.frame_email, textvariable=self.sender_email_var, width=40).grid(row=0, column=1, sticky="w", pady=5, padx=5)

        ttk.Label(self.frame_email, text="App Password:", font=("Helvetica", 10, "bold")).grid(row=1, column=0, sticky="w", pady=5, padx=5)
        self.sender_pass_var = tk.StringVar()
        ttk.Entry(self.frame_email, textvariable=self.sender_pass_var, show="*", width=40).grid(row=1, column=1, sticky="w", pady=5, padx=5)
        ttk.Label(self.frame_email, text="*Note: Use a Google App Password, not your regular login password.", font=("Helvetica", 8, "italic"), foreground="grey").grid(row=2, column=1, sticky="w", padx=5)

        self.frame_wa = ttk.Frame(self.comm_notebook, padding=15)
        self.comm_notebook.add(self.frame_wa, text="💬 WhatsApp (Meta WACA)")

        ttk.Label(self.frame_wa, text="Phone Number ID:", font=("Helvetica", 10, "bold")).grid(row=0, column=0, sticky="w", pady=5, padx=5)
        self.wa_phone_id_var = tk.StringVar()
        ttk.Entry(self.frame_wa, textvariable=self.wa_phone_id_var, width=45).grid(row=0, column=1, sticky="w", pady=5, padx=5)

        ttk.Label(self.frame_wa, text="Access Token:", font=("Helvetica", 10, "bold")).grid(row=1, column=0, sticky="w", pady=5, padx=5)
        self.wa_token_var = tk.StringVar()
        ttk.Entry(self.frame_wa, textvariable=self.wa_token_var, show="*", width=45).grid(row=1, column=1, sticky="w", pady=5, padx=5)
        
        ttk.Label(self.frame_wa, text="Template Name:", font=("Helvetica", 10, "bold")).grid(row=2, column=0, sticky="w", pady=5, padx=5)
        self.wa_template_var = tk.StringVar(value="lvlx_health_report")
        ttk.Entry(self.frame_wa, textvariable=self.wa_template_var, width=45).grid(row=2, column=1, sticky="w", pady=5, padx=5)
        ttk.Label(self.frame_wa, text="*Find credentials in your Meta Developers Dashboard -> WhatsApp -> API Setup.", font=("Helvetica", 8, "italic"), foreground="grey").grid(row=3, column=1, sticky="w", padx=5)

        action_frame = ttk.LabelFrame(self.tab_comm, text=" Execution Engine ", padding=15)
        action_frame.pack(fill="x", padx=20, pady=10)

        info_text = "Select your communication method above. This engine will scan the 'Master Report' folder for the specific date provided below, match the ERP to 'parents_data.csv', and automatically dispatch the reports."
        ttk.Label(action_frame, text=info_text, wraplength=1000).pack(anchor="w", pady=(0, 10))

        date_frame = ttk.Frame(action_frame)
        date_frame.pack(fill="x", pady=(0, 15))
        ttk.Label(date_frame, text="Target Date:", font=("Helvetica", 10, "bold")).pack(side="left")
        
        self.target_date_var = tk.StringVar(value=datetime.now().strftime("%d-%m-%Y"))
        
        if HAS_TKCALENDAR:
            self.date_picker = DateEntry(date_frame, textvariable=self.target_date_var, width=15, 
                                         background='#2980b9', foreground='white', borderwidth=2, 
                                         date_pattern='dd-mm-yyyy', font=("Helvetica", 10))
            self.date_picker.pack(side="left", padx=10)
        else:
            ttk.Entry(date_frame, textvariable=self.target_date_var, width=15).pack(side="left", padx=10)
            ttk.Label(date_frame, text="*(Tip: Install 'tkcalendar' for a visual date picker)*", font=("Helvetica", 8, "italic"), foreground="grey").pack(side="left")

        self.btn_send_reports = tk.Button(action_frame, text="🚀 Send Master Reports for Target Date", font=("Helvetica", 11, "bold"), command=self.start_distribution_thread, **self.get_btn_style("#8e44ad"))
        self.btn_send_reports.pack(anchor="w", ipadx=15, ipady=5)

        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(action_frame, variable=self.progress_var, maximum=100, length=400)
        self.progress_bar.pack(fill="x", pady=(15, 5))
        
        self.progress_lbl = ttk.Label(action_frame, text="Ready to send.", font=("Helvetica", 9))
        self.progress_lbl.pack(anchor="e")

        log_frame = ttk.LabelFrame(self.tab_comm, text=" Real-Time Status Log ", padding=10)
        log_frame.pack(fill="both", expand=True, padx=20, pady=(10, 20))

        self.comm_log = tk.Text(log_frame, state="disabled", font=("Courier", 9), bg="#1e1e1e", fg="#00ff00", wrap="word")
        self.comm_log.pack(fill="both", expand=True)

    def log_to_console(self, msg):
        self.root.after(0, self._log_gui, msg)
        if getattr(self, 'current_log_file', None):
            try:
                with open(self.current_log_file, 'a', encoding='utf-8') as f:
                    f.write(f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] {msg}\n")
            except Exception: pass 

    def _log_gui(self, msg):
        self.comm_log.config(state="normal")
        self.comm_log.insert("end", f"> {msg}\n")
        self.comm_log.see("end")
        self.comm_log.config(state="disabled")

    def update_progress(self, val, max_val, text):
        self.root.after(0, self._update_progress_gui, val, max_val, text)

    def _update_progress_gui(self, val, max_val, text):
        self.progress_bar["maximum"] = max_val
        self.progress_var.set(val)
        self.progress_lbl.config(text=text)

    def start_distribution_thread(self):
        target_date = self.target_date_var.get().strip()
        if not target_date: return messagebox.showwarning("Missing Date", "Please specify a Target Date (DD-MM-YYYY).")
        current_tab_index = self.comm_notebook.index(self.comm_notebook.select())
        
        if current_tab_index == 0:
            email, pwd = self.sender_email_var.get().strip(), self.sender_pass_var.get().strip()
            if not email or not pwd: return messagebox.showwarning("Missing Credentials", "Please enter your Sender Email and App Password.")
            if messagebox.askyesno("Confirm", f"Are you sure you want to begin sending EMAILS for {target_date}?\n\nThis will begin delivering attachments immediately."):
                self.btn_send_reports.config(state="disabled")
                self.comm_log.config(state="normal"); self.comm_log.delete("1.0", "end"); self.comm_log.config(state="disabled")
                threading.Thread(target=self.process_emails, args=(email, pwd, target_date), daemon=True).start()
        else:
            if not HAS_REQUESTS: return messagebox.showerror("Missing Library", "The 'requests' library is required for WhatsApp API.\n\nRun:\nuv add requests")
            phone_id, token, template_name = self.wa_phone_id_var.get().strip(), self.wa_token_var.get().strip(), self.wa_template_var.get().strip()
            if not phone_id or not token or not template_name: return messagebox.showwarning("Missing Credentials", "Please enter your Phone Number ID, Access Token, and Template Name.")
            if messagebox.askyesno("Confirm", f"Are you sure you want to begin sending via Meta WhatsApp API for {target_date}?\n\nThis will trigger live API requests immediately."):
                self.btn_send_reports.config(state="disabled")
                self.comm_log.config(state="normal"); self.comm_log.delete("1.0", "end"); self.comm_log.config(state="disabled")
                threading.Thread(target=self.process_whatsapp, args=(phone_id, token, template_name, target_date), daemon=True).start()

    def process_emails(self, sender_email, sender_pwd, target_date):
        db_filename = os.path.basename(self.db_name)
        school_name = os.path.splitext(db_filename)[0]
        school_dir = os.path.join(os.getcwd(), "Reports", school_name)
        
        os.makedirs(school_dir, exist_ok=True)
        logs_dir = os.path.join(school_dir, "logs")
        os.makedirs(logs_dir, exist_ok=True)
        self.current_log_file = os.path.join(logs_dir, "email_distribution_log.txt")
        
        self.log_to_console(f"INITIATING MASS EMAIL ENGINE FOR {target_date}...")
        
        daily_master_dir = os.path.join(school_dir, "Master Report", target_date)
        contact_csv = os.path.join(school_dir, "parents_data", "parents_data.csv")

        if not os.path.exists(daily_master_dir):
            self.log_to_console(f"ERROR: No Master Reports found for '{target_date}'. Have you built them yet?")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        if not os.path.exists(contact_csv):
            self.log_to_console("ERROR: 'parents_data.csv' not found. Please ensure it is placed in the 'parents_data' folder.")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        self.log_to_console("Parsing parents_data.csv...")
        contact_dict = {}
        try:
            with open(contact_csv, 'r', encoding='utf-8-sig') as f:
                reader = csv.DictReader(f)
                for row in reader:
                    erp_col = next((k for k in row.keys() if k and k.strip().upper() in ['ID', 'ERP']), None)
                    email_col = next((k for k in row.keys() if k and 'MAIL' in k.strip().upper()), None)
                    
                    if erp_col and email_col:
                        erp = str(row[erp_col]).strip()
                        mail = str(row[email_col]).strip()
                        if erp and mail:
                            contact_dict[erp] = mail
                            
            self.log_to_console(f"Successfully loaded {len(contact_dict)} email contacts.")
        except Exception as e:
            self.log_to_console(f"ERROR reading contacts CSV: {e}")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        pdf_files = [f for f in os.listdir(daily_master_dir) if f.startswith("LVLX_Master_Report_") and f.endswith(".pdf")]
        total_files = len(pdf_files)
        
        if total_files == 0:
            self.log_to_console(f"ERROR: Date folder '{target_date}' exists but contains no Master Reports.")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        self.log_to_console(f"Connecting to SMTP server (smtp.gmail.com)...")
        try:
            server = smtplib.SMTP("smtp.gmail.com", 587)
            server.starttls()
            server.login(sender_email, sender_pwd)
            self.log_to_console("SMTP Login Successful. Commencing distribution...")
        except Exception as e:
            self.log_to_console(f"SMTP LOGIN FAILED: {e}")
            self.log_to_console("Ensure you are using a valid Google App Password.")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        success_count = 0
        skip_count = 0
        
        for idx, filename in enumerate(pdf_files, start=1):
            erp = filename.replace("LVLX_Master_Report_", "").replace(".pdf", "")
            self.update_progress(idx, total_files, f"Processing {idx}/{total_files} - ERP: {erp}")
            
            target_email = contact_dict.get(erp)
            if not target_email:
                self.log_to_console(f"[SKIPPED] ERP {erp} - No matching email found.")
                skip_count += 1
                continue
                
            file_path = os.path.join(daily_master_dir, filename)
            files_to_attach = [file_path]

            try:
                msg = MIMEMultipart()
                msg['From'] = sender_email
                msg['To'] = target_email
                msg['Subject'] = "Your Child’s Health Snapshot is Ready – Discover What’s Next"
                
                body = f"""<html>
                <body>
                    <p>Dear Parent,</p>
                    <p>We’re pleased to share your child’s Health & Body Composition Report.</p>
                    <p>It offers key insights into their growth, fitness, and overall well-being.</p>
                    <p>For deeper analysis and personalized guidance, explore our premium program designed to support your child’s health journey.</p>
                    <p>Warm regards,<br>{self.display_school_name} & LVL X Junior<br>+91 9819300066</p>
                    <br>
                    <img src="cid:lvlx_logo" alt="LVL X Logo" width="180">
                </body>
                </html>"""
                
                msg.attach(MIMEText(body, 'html'))

                logo_path = os.path.join(os.getcwd(), "template", "lvlx_logo.png")
                if os.path.exists(logo_path):
                    with open(logo_path, "rb") as img_file:
                        logo_img = MIMEImage(img_file.read())
                        logo_img.add_header('Content-ID', '<lvlx_logo>')
                        logo_img.add_header('Content-Disposition', 'inline')
                        msg.attach(logo_img)

                for f_path in files_to_attach:
                    with open(f_path, "rb") as f:
                        part = MIMEApplication(f.read(), Name=os.path.basename(f_path))
                    part['Content-Disposition'] = f'attachment; filename="{os.path.basename(f_path)}"'
                    msg.attach(part)

                server.send_message(msg)
                self.log_to_console(f"[SENT] ERP {erp} -> {target_email}")
                success_count += 1
                time.sleep(1)
            except Exception as e:
                self.log_to_console(f"[FAILED] ERP {erp} -> {target_email}. Error: {e}")
                skip_count += 1

        try: server.quit() 
        except: pass

        self.log_to_console("="*40)
        self.log_to_console(f"DISTRIBUTION COMPLETE. Sent: {success_count}, Skipped: {skip_count}")
        self.update_progress(total_files, total_files, f"Finished! Sent: {success_count} | Skipped: {skip_count}")
        self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
        self.current_log_file = None 

    def process_whatsapp(self, phone_id, token, template_name, target_date):
        db_filename = os.path.basename(self.db_name)
        school_name = os.path.splitext(db_filename)[0]
        school_dir = os.path.join(os.getcwd(), "Reports", school_name)
        
        os.makedirs(school_dir, exist_ok=True)
        logs_dir = os.path.join(school_dir, "logs")
        os.makedirs(logs_dir, exist_ok=True)
        self.current_log_file = os.path.join(logs_dir, "whatsapp_distribution_log.txt")

        self.log_to_console(f"INITIATING META WACA ENGINE FOR {target_date}...")
        
        daily_master_dir = os.path.join(school_dir, "Master Report", target_date)
        contact_csv = os.path.join(school_dir, "parents_data", "parents_data.csv")

        if not os.path.exists(daily_master_dir):
            self.log_to_console(f"ERROR: No Master Reports found for '{target_date}'. Have you built them yet?")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        if not os.path.exists(contact_csv):
            self.log_to_console("ERROR: 'parents_data.csv' not found. Please ensure it is placed in the 'parents_data' folder.")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        self.log_to_console("Parsing parents_data.csv for Phone Numbers...")
        contact_dict = {}
        try:
            with open(contact_csv, 'r', encoding='utf-8-sig') as f:
                reader = csv.DictReader(f)
                for row in reader:
                    erp_col = next((k for k in row.keys() if k and k.strip().upper() in ['ID', 'ERP']), None)
                    phone_col = next((k for k in row.keys() if k and 'NUM' in k.strip().upper()), None)
                    
                    if erp_col and phone_col:
                        erp = str(row[erp_col]).strip()
                        phone = str(row[phone_col]).strip()
                        phone_clean = ''.join(filter(str.isdigit, phone))
                        if erp and phone_clean:
                            contact_dict[erp] = phone_clean
                            
            self.log_to_console(f"Successfully loaded {len(contact_dict)} phone numbers.")
        except Exception as e:
            self.log_to_console(f"ERROR reading contacts CSV: {e}")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        pdf_files = [f for f in os.listdir(daily_master_dir) if f.startswith("LVLX_Master_Report_") and f.endswith(".pdf")]
        total_files = len(pdf_files)
        
        if total_files == 0:
            self.log_to_console(f"ERROR: Date folder '{target_date}' exists but contains no Master Reports.")
            self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
            self.current_log_file = None
            return

        success_count = 0
        skip_count = 0
        
        upload_url = f"https://graph.facebook.com/v19.0/{phone_id}/media"
        send_url = f"https://graph.facebook.com/v19.0/{phone_id}/messages"
        headers = {"Authorization": f"Bearer {token}"}
        
        for idx, filename in enumerate(pdf_files, start=1):
            erp = filename.replace("LVLX_Master_Report_", "").replace(".pdf", "")
            self.update_progress(idx, total_files, f"Processing {idx}/{total_files} - ERP: {erp}")
            
            target_phone = contact_dict.get(erp)
            if not target_phone:
                self.log_to_console(f"[SKIPPED] ERP {erp} - No phone number found.")
                skip_count += 1
                continue
                
            file_path = os.path.join(daily_master_dir, filename)
            files_to_attach = [file_path]

            try:
                self.log_to_console(f"[PREPARING] Uploading final PDF for ERP {erp} to Meta Servers...")
                
                for f_path in files_to_attach:
                    file_name = os.path.basename(f_path)
                    
                    with open(f_path, "rb") as f:
                        files = {"file": (file_name, f, "application/pdf")}
                        data = {"messaging_product": "whatsapp", "type": "application/pdf"}
                        
                        resp = requests.post(upload_url, headers=headers, data=data, files=files)
                        resp_json = resp.json()
                        
                        if "id" not in resp_json:
                            raise Exception(f"Failed to upload {file_name}: {resp_json}")
                            
                        media_id = resp_json["id"]

                    media_payload = {
                        "messaging_product": "whatsapp",
                        "recipient_type": "individual",
                        "to": target_phone,
                        "type": "template",
                        "template": {
                            "name": template_name,
                            "language": {"code": "en"},
                            "components": [
                                {
                                    "type": "header",
                                    "parameters": [
                                        {
                                            "type": "document",
                                            "document": {
                                                "id": media_id,
                                                "filename": file_name
                                            }
                                        }
                                    ]
                                }
                            ]
                        }
                    }
                    
                    send_resp = requests.post(send_url, headers=headers, json=media_payload)
                    if send_resp.status_code not in [200, 201]:
                        raise Exception(f"Failed to deliver {file_name}: {send_resp.text}")
                        
                    time.sleep(1) 

                self.log_to_console(f"[SENT] ERP {erp} -> WhatsApp: {target_phone} ({len(files_to_attach)} file)")
                success_count += 1
                
            except Exception as e:
                self.log_to_console(f"[FAILED] ERP {erp} -> {target_phone}. Error: {e}")
                skip_count += 1

        self.log_to_console("="*40)
        self.log_to_console(f"WHATSAPP DISTRIBUTION COMPLETE. Sent: {success_count}, Skipped: {skip_count}")
        self.update_progress(total_files, total_files, f"Finished! Sent: {success_count} | Skipped: {skip_count}")
        self.root.after(0, lambda: self.btn_send_reports.config(state="normal"))
        self.current_log_file = None 

    def setup_analytics_tab(self):
        filter_frame = ttk.Frame(self.tab_analytics)
        filter_frame.pack(fill="x", padx=10, pady=(10, 0))
        
        ttk.Label(filter_frame, text="View Metric: ", font=("Helvetica", 11, "bold")).pack(side="left")
        
        self.metric_var = tk.StringVar(value="Overview")
        cb = ttk.Combobox(filter_frame, textvariable=self.metric_var, values=["Overview", "Demographics", "BMI", "Height", "Weight"], state="readonly", width=25)
        cb.pack(side="left", padx=10)
        cb.bind("<<ComboboxSelected>>", lambda e: self.update_analytics_view())

        tk.Button(filter_frame, text="📸 Export All Charts to Folder", font=("Helvetica", 9, "bold"), command=self.export_charts, **self.get_btn_style("#8e44ad")).pack(side="right", padx=10)

        self.kpi_frame = ttk.LabelFrame(self.tab_analytics, text=" Key Performance Indicators (KPIs) ", padding=10)
        self.kpi_frame.pack(fill="x", pady=(10, 10), padx=10)

        self.chart_frame = ttk.Frame(self.tab_analytics)
        self.chart_frame.pack(fill="both", expand=True, padx=10, pady=10)

    def export_charts(self):
        if getattr(self, 'total_students', 0) == 0: return messagebox.showwarning("No Data", "No data to export!")
        if not HAS_PILLOW: return messagebox.showerror("Missing Library", "The 'Pillow' library is required to take screenshots.\n\nRun:\nuv add pillow")

        db_filename = os.path.basename(self.db_name)
        school_name = os.path.splitext(db_filename)[0]
        
        export_dir = os.path.join(os.getcwd(), "Reports", school_name, "exports", "charts_data")
        os.makedirs(export_dir, exist_ok=True)
        
        safe_school_name = self.display_school_name.replace(" ", "_")
        current_metric = self.metric_var.get()
        metrics = ["Overview", "Demographics", "BMI", "Height", "Weight"]

        self.root.config(cursor="wait"); self.root.update()

        try:
            for metric in metrics:
                self.metric_var.set(metric)
                self.update_analytics_view()
                self.root.update_idletasks()
                self.root.update()
                time.sleep(0.3) 
                
                x1, y1 = self.kpi_frame.winfo_rootx(), self.kpi_frame.winfo_rooty()
                x2, y2 = self.chart_frame.winfo_rootx() + self.chart_frame.winfo_width(), self.chart_frame.winfo_rooty() + self.chart_frame.winfo_height()
                
                img = ImageGrab.grab(bbox=(x1, y1, x2, y2))
                img_path = os.path.join(export_dir, f"{safe_school_name}_{metric}_Dashboard.png")
                img.save(img_path, format='PNG')
                
            messagebox.showinfo("Success", f"All charts exported successfully as PNGs!\n\nLocation:\n{export_dir}")
        except Exception as e: messagebox.showerror("Export Error", f"Failed to export charts:\n{str(e)}")
        finally:
            self.metric_var.set(current_metric); self.update_analytics_view(); self.root.config(cursor="")

    def update_analytics_view(self):
        for w in self.kpi_frame.winfo_children(): w.destroy()
        for w in self.chart_frame.winfo_children(): w.destroy()

        if getattr(self, 'total_students', 0) == 0:
            ttk.Label(self.kpi_frame, text="No data available.").pack()
            return

        metric = self.metric_var.get()

        if metric == "Overview":
            total_boys = sum(1 for r in self.all_student_data if r[4] == 'M')
            total_girls = sum(1 for r in self.all_student_data if r[4] == 'F')

            tk.Label(self.kpi_frame, text=f"Total Students\n{self.total_students}", font=("Helvetica", 12, "bold"), bg="#3498db", fg="white", pady=10).pack(side="left", padx=5, expand=True, fill="both")
            tk.Label(self.kpi_frame, text=f"Total Boys\n{total_boys}", font=("Helvetica", 12, "bold"), bg="#2980b9", fg="white", pady=10).pack(side="left", padx=5, expand=True, fill="both")
            tk.Label(self.kpi_frame, text=f"Total Girls\n{total_girls}", font=("Helvetica", 12, "bold"), bg="#8e44ad", fg="white", pady=10).pack(side="left", padx=5, expand=True, fill="both")

        elif metric == "Demographics":
            tk.Label(self.kpi_frame, text=f"Total Students\n{self.total_students}", font=("Helvetica", 12, "bold"), bg="#3498db", fg="white", pady=10).pack(side="left", padx=5, expand=True, fill="both")
            for cls, stats in self.class_stats.items():
                tk.Label(self.kpi_frame, text=f"Class {cls}\nTotal: {stats['Total']} (M: {stats['M']} | F: {stats['F']})", font=("Helvetica", 10), bg="#ecf0f1", pady=10, relief="solid", borderwidth=1).pack(side="left", padx=5, expand=True, fill="both")
        
        else:
            total_wr = sum(c['Within Range']['M'] + c['Within Range']['F'] for c in self.cat_stats[metric].values())
            total_b = sum(c['Borderline']['M'] + c['Borderline']['F'] for c in self.cat_stats[metric].values())
            total_na = sum(c['Needs Attention']['M'] + c['Needs Attention']['F'] for c in self.cat_stats[metric].values())

            tk.Label(self.kpi_frame, text=f"Within Range\n{total_wr}", font=("Helvetica", 12, "bold"), bg="#2ecc71", fg="white", pady=10).pack(side="left", padx=5, expand=True, fill="both")
            tk.Label(self.kpi_frame, text=f"Borderline\n{total_b}", font=("Helvetica", 12, "bold"), bg="#f1c40f", fg="black", pady=10).pack(side="left", padx=5, expand=True, fill="both")
            tk.Label(self.kpi_frame, text=f"Needs Attention\n{total_na}", font=("Helvetica", 12, "bold"), bg="#e67e22", fg="white", pady=10).pack(side="left", padx=5, expand=True, fill="both")

            for cls, stats in self.cat_stats[metric].items():
                cls_wr = stats['Within Range']['M'] + stats['Within Range']['F']
                cls_b = stats['Borderline']['M'] + stats['Borderline']['F']
                cls_na = stats['Needs Attention']['M'] + stats['Needs Attention']['F']
                tk.Label(self.kpi_frame, text=f"Class {cls}\nIn Range: {cls_wr} | Border: {cls_b} | Attn: {cls_na}", font=("Helvetica", 10), bg="#ecf0f1", pady=10, relief="solid", borderwidth=1).pack(side="left", padx=5, expand=True, fill="both")

        self.draw_chart(metric)

    def _add_tooltips(self, ax, rects_list, fs=9):
        for rects in rects_list: 
            ax.bar_label(rects, label_type='center', color='white', padding=3, fmt='%d', fontweight='bold', fontsize=fs)

    def draw_chart(self, metric):
        if self.is_mac:
            c_fig_size = (7.5, 4.5) 
            c_dpi = 80
            c_title_fs = 12
            c_label_fs = 8
            c_pie_fs = 7
            c_pie_title_fs = 10
            c_hspace = 0.4
            c_wspace = 0.3
        else:
            c_fig_size = (10, 6)
            c_dpi = 100
            c_title_fs = 14
            c_label_fs = 9
            c_pie_fs = 8
            c_pie_title_fs = 11
            c_hspace = 0.3
            c_wspace = 0.2

        fig = plt.figure(figsize=c_fig_size, dpi=c_dpi)

        if metric == "Overview":
            axes = [fig.add_subplot(221), fig.add_subplot(222), fig.add_subplot(223), fig.add_subplot(224)]
            titles = ['Boys Height', 'Girls Height', 'Boys Weight', 'Girls Weight']
            labels, colors = ['Within Range', 'Borderline', 'Needs Attention'], ['#2ecc71', '#f1c40f', '#e67e22']

            for ax, title in zip(axes, titles):
                data = self.overview_stats_4q[title]
                sizes = [data['Within Range'], data['Borderline'], data['Needs Attention']]
                sizes_f = [s for s in sizes if s > 0]
                labels_f = [l for s, l in zip(sizes, labels) if s > 0]
                colors_f = [c for s, c in zip(sizes, colors) if s > 0]

                if sizes_f:
                    wedges, texts, autotexts = ax.pie(
                        sizes_f, 
                        labels=labels_f, 
                        colors=colors_f, 
                        autopct=lambda p: f"{p:.1f}%\n({int(round(p * sum(sizes_f) / 100.0))})", 
                        startangle=90, 
                        textprops={'fontweight': 'bold', 'fontsize': c_pie_fs}
                    )
                    for i, autotext in enumerate(autotexts): 
                        autotext.set_color('black' if labels_f[i] == 'Borderline' else 'white')
                else: 
                    ax.text(0.5, 0.5, 'No Data', ha='center', va='center', fontweight='bold', color="#7f8c8d", fontsize=c_label_fs)
                    
                ax.set_title(title, fontweight='bold', pad=8, fontsize=c_pie_title_fs)
                ax.axis('equal')

            fig.suptitle('School Health Overview: Height & Weight by Gender', fontweight='bold', fontsize=c_title_fs)
            
            if self.is_mac:
                fig.subplots_adjust(top=0.85, bottom=0.05, hspace=c_hspace, wspace=c_wspace)
            else:
                fig.subplots_adjust(top=0.90, bottom=0.05, hspace=c_hspace)

        elif metric == "Demographics":
            ax = fig.add_subplot(111)
            classes, x, width = list(self.class_stats.keys()), np.arange(len(self.class_stats.keys())), 0.4
            boys = [self.class_stats[c]['M'] for c in classes]
            girls = [self.class_stats[c]['F'] for c in classes]

            r1 = ax.bar(x, boys, width, label='Boys', color='#3498db')
            r2 = ax.bar(x, girls, width, bottom=boys, label='Girls', color='#fd79a8')

            ax.set_title('Student Demographics by Class', fontweight='bold', pad=15, fontsize=c_title_fs)
            self._add_tooltips(ax, [r1, r2], c_label_fs)
            ax.set_xticks(x); ax.set_xticklabels(classes, fontweight='bold', fontsize=c_label_fs)
            ax.set_ylabel('Number of Students', fontweight='bold', fontsize=c_label_fs)
            ax.legend(fontsize=c_label_fs)
            fig.tight_layout()

        else:
            ax = fig.add_subplot(111)
            classes, x, width = list(self.cat_stats[metric].keys()), np.arange(len(self.cat_stats[metric].keys())), 0.25
            wr_m, wr_f = [self.cat_stats[metric][c]['Within Range']['M'] for c in classes], [self.cat_stats[metric][c]['Within Range']['F'] for c in classes]
            b_m, b_f = [self.cat_stats[metric][c]['Borderline']['M'] for c in classes], [self.cat_stats[metric][c]['Borderline']['F'] for c in classes]
            na_m, na_f = [self.cat_stats[metric][c]['Needs Attention']['M'] for c in classes], [self.cat_stats[metric][c]['Needs Attention']['F'] for c in classes]

            r1_m = ax.bar(x - width, wr_m, width, color='#27ae60', label='Boys (Within Range)')
            r1_f = ax.bar(x - width, wr_f, width, bottom=wr_m, color='#2ecc71', label='Girls (Within Range)')
            r2_m = ax.bar(x, b_m, width, color='#f39c12', label='Boys (Borderline)')
            r2_f = ax.bar(x, b_f, width, bottom=b_m, color='#f1c40f', label='Girls (Borderline)')
            r3_m = ax.bar(x + width, na_m, width, color='#d35400', label='Boys (Needs Attention)')
            r3_f = ax.bar(x + width, na_f, width, bottom=na_m, color='#e67e22', label='Girls (Needs Attention)')

            ax.set_title(f'{metric} Segregation by Class', fontweight='bold', pad=15, fontsize=c_title_fs)
            self._add_tooltips(ax, [r1_m, r1_f, r2_m, r2_f, r3_m, r3_f], c_label_fs)
            ax.set_xticks(x); ax.set_xticklabels(classes, fontweight='bold', fontsize=c_label_fs)
            ax.legend(loc='upper left', bbox_to_anchor=(1, 1), fontsize=c_label_fs)
            ax.set_ylabel('Number of Students', fontweight='bold', fontsize=c_label_fs)
            
            fig.tight_layout()
            if self.is_mac:
                fig.subplots_adjust(right=0.70) 

        canvas = FigureCanvasTkAgg(fig, master=self.chart_frame)
        canvas.draw()
        canvas.get_tk_widget().pack(fill="both", expand=True)

    def refresh_all_data(self):
        if not os.path.exists(self.db_name): return

        try:
            with sqlite3.connect(self.db_name) as conn:
                cursor = conn.cursor()

                cursor.execute("SELECT COUNT(*) FROM growth_reports")
                self.total_students = cursor.fetchone()[0]

                cursor.execute("SELECT class_grade, gender, COUNT(*) FROM growth_reports GROUP BY class_grade, gender ORDER BY class_grade")
                self.class_stats = {}
                for cls, gender, count in cursor.fetchall():
                    cls = cls or "Unassigned"
                    if cls not in self.class_stats:
                        self.class_stats[cls] = {'M': 0, 'F': 0, 'Total': 0}
                    if gender in ['M', 'F']:
                        self.class_stats[cls][gender] += count
                    self.class_stats[cls]['Total'] += count

                self.cat_stats = {m: {c: {'Within Range': {'M': 0, 'F': 0}, 'Borderline': {'M': 0, 'F': 0}, 'Needs Attention': {'M': 0, 'F': 0}} for c in self.class_stats} for m in ['BMI', 'Height', 'Weight']}
                
                self.overview_stats_4q = {
                    'Boys Height': {'Within Range': 0, 'Borderline': 0, 'Needs Attention': 0},
                    'Girls Height': {'Within Range': 0, 'Borderline': 0, 'Needs Attention': 0},
                    'Boys Weight': {'Within Range': 0, 'Borderline': 0, 'Needs Attention': 0},
                    'Girls Weight': {'Within Range': 0, 'Borderline': 0, 'Needs Attention': 0}
                }

                cursor.execute("SELECT class_grade, age, gender, height, weight, bmi FROM growth_reports WHERE class_grade != '' AND age != ''")
                for cls, age_str, gender, h, w, bmi in cursor.fetchall():
                    try: age = int(str(age_str).split('/')[0])
                    except ValueError: continue
                    
                    if cls not in self.cat_stats['BMI']: continue
                    
                    if bmi: 
                        b_cat = categorize_metric('BMI', float(bmi), age, gender)
                        self.cat_stats['BMI'][cls][b_cat][gender] += 1
                        
                    if h: 
                        h_cat = categorize_metric('Height', float(h), age, gender)
                        self.cat_stats['Height'][cls][h_cat][gender] += 1
                        if gender == 'M': self.overview_stats_4q['Boys Height'][h_cat] += 1
                        elif gender == 'F': self.overview_stats_4q['Girls Height'][h_cat] += 1
                        
                    if w: 
                        w_cat = categorize_metric('Weight', float(w), age, gender)
                        self.cat_stats['Weight'][cls][w_cat][gender] += 1
                        if gender == 'M': self.overview_stats_4q['Boys Weight'][w_cat] += 1
                        elif gender == 'F': self.overview_stats_4q['Girls Weight'][w_cat] += 1

                cursor.execute("SELECT * FROM growth_reports ORDER BY id DESC")
                self.all_student_data = cursor.fetchall()

            self.update_analytics_view()
            
            self.tree.delete(*self.tree.get_children())
            for row in self.all_student_data: 
                self.tree.insert("", "end", values=row)

        except sqlite3.Error as e:
            messagebox.showerror("Database Error", str(e))

    def on_closing(self):
        try:
            plt.close('all')
            self.root.quit()
            self.root.destroy()
        except Exception: pass
        finally: os._exit(0)

if __name__ == "__main__":
    root = tk.Tk()
    app = LVLXCommandCenter(root)
    root.mainloop()